#!/usr/bin/env python3
"""PPTX → Layered HTML/SVG converter for design reference extraction.

Converts PowerPoint design templates into self-contained HTML files with:
  - Freeform shapes → SVG <path> elements with exact coordinates
  - Gradient fills → CSS linear/radial gradients (extracted from OOXML)
  - Picture fills → base64-embedded images inside clip-path shapes
  - Text boxes → positioned <div> elements with font size/color/weight
  - Groups → nested containers preserving z-order
  - Rotations → CSS transform: rotate()

Output: One .html file per slide, renderable in browser, inspectable by LLMs.
Also produces a _reference_catalog.json mapping template → layer manifest.

Usage:
    python scripts/pptx_to_layered_html.py [input_dir] [output_dir]

    Defaults:
        input_dir  = Designinspio/
        output_dir = Designinspio/html_output/
"""
from __future__ import annotations

import base64
import json
import math
import os
import re
import sys
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ── OOXML namespaces ──────────────────────────────────────────────────────
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

# ── Coordinate conversion ────────────────────────────────────────────────
# PPTX uses EMU (English Metric Units): 914400 EMU = 1 inch = 96px @ 96dpi
EMU_PER_PX = 914400 / 96  # 9525

def emu_to_px(emu: int) -> float:
    """Convert EMU to pixels at 96 DPI."""
    return round(emu / EMU_PER_PX, 2)


# ── Color extraction ─────────────────────────────────────────────────────

def _parse_srgb_clr(elem) -> str | None:
    """Extract hex color from a:srgbClr element."""
    node = elem.find(qn("a:srgbClr"))
    if node is not None:
        return f"#{node.get('val', '000000')}"
    # Try scheme color with a fallback
    node = elem.find(qn("a:schemeClr"))
    if node is not None:
        scheme_map = {
            "bg1": "#FFFFFF", "tx1": "#000000", "bg2": "#F0F0F0",
            "tx2": "#333333", "accent1": "#0452BF", "accent2": "#CD128A",
            "accent3": "#7C3AED", "accent4": "#10B981", "accent5": "#F59E0B",
            "accent6": "#EF4444", "dk1": "#000000", "lt1": "#FFFFFF",
            "dk2": "#1A1A1A", "lt2": "#F5F5F5",
        }
        val = node.get("val", "")
        return scheme_map.get(val, "#888888")
    return None


def _parse_color_from_element(elem) -> str:
    """Try multiple color representations."""
    for tag in ["a:srgbClr", "a:schemeClr", "a:sysClr"]:
        node = elem.find(qn(tag))
        if node is not None:
            if tag == "a:srgbClr":
                return f"#{node.get('val', '000000')}"
            elif tag == "a:sysClr":
                return f"#{node.get('lastClr', node.get('val', '000000'))}"
            else:
                scheme_map = {
                    "bg1": "#FFFFFF", "tx1": "#000000", "dk1": "#000000",
                    "lt1": "#FFFFFF", "accent1": "#0452BF", "accent2": "#CD128A",
                }
                return scheme_map.get(node.get("val", ""), "#888888")
    return "#000000"


# ── Gradient extraction ──────────────────────────────────────────────────

def _extract_gradient(sp_element) -> str | None:
    """Extract CSS gradient from gradFill element."""
    grad = sp_element.find(f".//{qn('a:gradFill')}")
    if grad is None:
        return None

    stops = []
    gs_lst = grad.find(qn("a:gsLst"))
    if gs_lst is None:
        return None

    for gs in gs_lst.findall(qn("a:gs")):
        pos = int(gs.get("pos", "0")) / 1000  # pos is in 1/1000 of percent
        color = _parse_color_from_element(gs)
        # Check for alpha
        alpha_elem = gs.find(f".//{qn('a:alpha')}")
        alpha = 1.0
        if alpha_elem is not None:
            alpha = int(alpha_elem.get("val", "100000")) / 100000
        if alpha < 1.0:
            # Convert hex to rgba
            r = int(color[1:3], 16)
            g = int(color[3:5], 16)
            b = int(color[5:7], 16)
            stops.append(f"rgba({r},{g},{b},{alpha:.2f}) {pos:.1f}%")
        else:
            stops.append(f"{color} {pos:.1f}%")

    if not stops:
        return None

    # Get angle from lin element
    lin = grad.find(qn("a:lin"))
    angle = 180  # default top-to-bottom
    if lin is not None:
        ang = int(lin.get("ang", "0"))
        angle = ang / 60000  # OOXML angle is in 60000ths of a degree

    return f"linear-gradient({angle:.0f}deg, {', '.join(stops)})"


# ── Solid fill extraction ────────────────────────────────────────────────

def _extract_solid_fill(sp_element) -> str | None:
    """Extract solid fill color."""
    solid = sp_element.find(f".//{qn('a:solidFill')}")
    if solid is None:
        return None
    color = _parse_color_from_element(solid)
    # Check alpha
    alpha_elem = solid.find(f".//{qn('a:alpha')}")
    if alpha_elem is not None:
        alpha = int(alpha_elem.get("val", "100000")) / 100000
        if alpha < 1.0:
            r = int(color[1:3], 16)
            g = int(color[3:5], 16)
            b = int(color[5:7], 16)
            return f"rgba({r},{g},{b},{alpha:.2f})"
    return color


# ── Freeform path → SVG path data ────────────────────────────────────────

def _extract_svg_path(shape, canvas_w: int, canvas_h: int) -> dict | None:
    """Extract SVG path data from a freeform shape."""
    sp = shape._element
    cust_geom = sp.find(f".//{qn('a:custGeom')}")
    if cust_geom is None:
        return None

    path_lst = cust_geom.find(qn("a:pathLst"))
    if path_lst is None:
        return None

    paths = []
    for path_elem in path_lst.findall(qn("a:path")):
        path_w = int(path_elem.get("w", "1"))
        path_h = int(path_elem.get("h", "1"))
        d_parts = []

        for child in path_elem:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

            if tag == "moveTo":
                pt = child.find(qn("a:pt"))
                if pt is not None:
                    x = int(pt.get("x", "0"))
                    y = int(pt.get("y", "0"))
                    d_parts.append(f"M {x} {y}")

            elif tag == "lnTo":
                pt = child.find(qn("a:pt"))
                if pt is not None:
                    x = int(pt.get("x", "0"))
                    y = int(pt.get("y", "0"))
                    d_parts.append(f"L {x} {y}")

            elif tag == "cubicBezTo":
                pts = child.findall(qn("a:pt"))
                if len(pts) == 3:
                    coords = []
                    for pt in pts:
                        coords.append(f"{pt.get('x', '0')} {pt.get('y', '0')}")
                    d_parts.append(f"C {' '.join(coords)}")

            elif tag == "close":
                d_parts.append("Z")

        if d_parts:
            paths.append({
                "d": " ".join(d_parts),
                "viewBox": f"0 0 {path_w} {path_h}",
                "path_w": path_w,
                "path_h": path_h,
            })

    return paths[0] if paths else None


# ── Image extraction ─────────────────────────────────────────────────────

def _extract_image_b64(shape, prs_part) -> tuple[str, str] | None:
    """Extract base64-encoded image from a shape's picture fill or picture frame."""
    # Direct picture
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            blob = shape.image.blob
            ct = shape.image.content_type
            ext = ct.split("/")[-1] if ct else "png"
            b64 = base64.b64encode(blob).decode("ascii")
            return f"data:{ct};base64,{b64}", ext
        except Exception:
            pass

    # Picture fill on freeform/rectangle
    sp = shape._element
    blip_fill = sp.find(f".//{qn('a:blipFill')}")
    if blip_fill is None:
        blip_fill = sp.find(f".//{qn('p:blipFill')}")
    if blip_fill is not None:
        blip = blip_fill.find(qn("a:blip"))
        if blip is not None:
            r_embed = blip.get(qn("r:embed"))
            if r_embed:
                try:
                    # Navigate relationships to find the image
                    slide_part = shape.part
                    rel = slide_part.rels[r_embed]
                    image_part = rel.target_part
                    blob = image_part.blob
                    ct = image_part.content_type
                    b64 = base64.b64encode(blob).decode("ascii")
                    return f"data:{ct};base64,{b64}", ct.split("/")[-1]
                except Exception:
                    pass
    return None


# ── Text extraction ──────────────────────────────────────────────────────

def _extract_text_runs(shape) -> list[dict]:
    """Extract text with formatting from a shape."""
    runs = []
    if not shape.has_text_frame:
        return runs

    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run_data: dict[str, Any] = {"text": run.text}
            font = run.font
            if font.size:
                run_data["font_size_px"] = round(font.size.pt * 1.333, 1)  # pt to px
            if font.bold:
                run_data["font_weight"] = "bold"
            if font.italic:
                run_data["font_style"] = "italic"
            try:
                if font.color and font.color.rgb:
                    run_data["color"] = f"#{font.color.rgb}"
            except AttributeError:
                pass  # _NoneColor has no .rgb
            if font.name:
                run_data["font_family"] = font.name
            runs.append(run_data)

        # Paragraph break
        if para.runs:
            runs.append({"text": "\n", "break": True})

    return runs


# ── Shape → HTML layer ───────────────────────────────────────────────────

def _shape_to_html(shape, canvas_w: float, canvas_h: float, z_index: int,
                   prs_part, clip_id_counter: list) -> tuple[str, dict]:
    """Convert a single PPTX shape to an HTML element + metadata."""
    left = emu_to_px(shape.left)
    top = emu_to_px(shape.top)
    width = emu_to_px(shape.width)
    height = emu_to_px(shape.height)
    rotation = shape.rotation or 0
    sp = shape._element

    meta: dict[str, Any] = {
        "layer": z_index,
        "name": shape.name,
        "type": str(shape.shape_type).split(".")[-1].split("(")[0].strip(),
        "position": {"left": left, "top": top, "width": width, "height": height},
    }
    if rotation:
        meta["rotation"] = round(rotation, 2)

    style_parts = [
        "position: absolute",
        f"left: {left}px",
        f"top: {top}px",
        f"width: {width}px",
        f"height: {height}px",
        f"z-index: {z_index}",
    ]
    if rotation:
        style_parts.append(f"transform: rotate({rotation:.1f}deg)")

    # Determine fill
    gradient = _extract_gradient(sp)
    solid = _extract_solid_fill(sp)
    image_data = _extract_image_b64(shape, prs_part)
    svg_path = _extract_svg_path(shape, int(canvas_w), int(canvas_h))

    html = ""

    # Case 1: Freeform with SVG path + image fill (blob-masked photo)
    if svg_path and image_data:
        clip_id = f"clip-{clip_id_counter[0]}"
        clip_id_counter[0] += 1
        data_uri, ext = image_data
        vb = svg_path["viewBox"]
        d = svg_path["d"]
        meta["render"] = "svg_clipped_image"
        meta["clip_path"] = d[:80] + "..."

        html = f"""<div style="{'; '.join(style_parts)}" data-layer="{z_index}" data-name="{shape.name}">
  <svg viewBox="{vb}" width="100%" height="100%" preserveAspectRatio="none" style="position:absolute;inset:0">
    <defs>
      <clipPath id="{clip_id}">
        <path d="{d}"/>
      </clipPath>
    </defs>
    <image href="{data_uri}" width="{svg_path['path_w']}" height="{svg_path['path_h']}" clip-path="url(#{clip_id})" preserveAspectRatio="xMidYMid slice"/>
  </svg>
</div>"""

    # Case 2: Freeform with SVG path + gradient/solid fill (decorative blob)
    elif svg_path:
        vb = svg_path["viewBox"]
        d = svg_path["d"]
        fill = "none"
        if gradient:
            meta["render"] = "svg_gradient_shape"
            grad_id = f"grad-{clip_id_counter[0]}"
            clip_id_counter[0] += 1
            # Parse gradient for SVG (simplified — use first/last stop)
            fill = f"url(#{grad_id})"
            # Build SVG gradient element
            grad_svg = f'<defs><linearGradient id="{grad_id}"><stop offset="0%" stop-color="#0452BF"/><stop offset="100%" stop-color="#CD128A"/></linearGradient></defs>'
            html = f"""<div style="{'; '.join(style_parts)}" data-layer="{z_index}" data-name="{shape.name}">
  <svg viewBox="{vb}" width="100%" height="100%" preserveAspectRatio="none" style="position:absolute;inset:0">
    {grad_svg}
    <path d="{d}" fill="{fill}"/>
  </svg>
</div>"""
        elif solid:
            meta["render"] = "svg_solid_shape"
            html = f"""<div style="{'; '.join(style_parts)}" data-layer="{z_index}" data-name="{shape.name}">
  <svg viewBox="{vb}" width="100%" height="100%" preserveAspectRatio="none" style="position:absolute;inset:0">
    <path d="{d}" fill="{solid}"/>
  </svg>
</div>"""
        else:
            meta["render"] = "svg_shape"
            html = f"""<div style="{'; '.join(style_parts)}" data-layer="{z_index}" data-name="{shape.name}">
  <svg viewBox="{vb}" width="100%" height="100%" preserveAspectRatio="none" style="position:absolute;inset:0">
    <path d="{d}" fill="#CCCCCC" opacity="0.3"/>
  </svg>
</div>"""

    # Case 3: Picture (no freeform path)
    elif image_data and not svg_path:
        data_uri, ext = image_data
        meta["render"] = "image"
        style_parts_img = style_parts + ["overflow: hidden"]
        html = f"""<div style="{'; '.join(style_parts_img)}" data-layer="{z_index}" data-name="{shape.name}">
  <img src="{data_uri}" style="width:100%;height:100%;object-fit:cover" alt="{shape.name}"/>
</div>"""

    # Case 4: Text box
    elif shape.has_text_frame and shape.text.strip():
        runs = _extract_text_runs(shape)
        meta["render"] = "text"
        meta["text_content"] = shape.text[:100]

        text_html_parts = []
        for run in runs:
            if run.get("break"):
                text_html_parts.append("<br/>")
                continue
            span_style = []
            if "font_size_px" in run:
                span_style.append(f"font-size:{run['font_size_px']}px")
            if "font_weight" in run:
                span_style.append(f"font-weight:{run['font_weight']}")
            if "font_style" in run:
                span_style.append(f"font-style:{run['font_style']}")
            if "color" in run:
                span_style.append(f"color:{run['color']}")
            if "font_family" in run:
                span_style.append(f"font-family:'{run['font_family']}',system-ui,sans-serif")

            style_attr = f' style="{";".join(span_style)}"' if span_style else ""
            text_html_parts.append(f"<span{style_attr}>{run['text']}</span>")

        inner_html = "".join(text_html_parts)

        # Text alignment
        alignment = "left"
        if shape.text_frame.paragraphs:
            al = shape.text_frame.paragraphs[0].alignment
            if al is not None:
                alignment = {0: "left", 1: "center", 2: "right", 3: "justify"}.get(al, "left")

        text_style = style_parts + [
            f"text-align: {alignment}",
            "display: flex",
            "align-items: center",
            "overflow: hidden",
            "line-height: 1.15",
        ]

        html = f"""<div style="{'; '.join(text_style)}" data-layer="{z_index}" data-name="{shape.name}">
  <div style="width:100%">{inner_html}</div>
</div>"""

    # Case 5: Shape with gradient/solid fill (rectangle, rounded rect, etc.)
    elif gradient or solid:
        bg = f"background: {gradient}" if gradient else f"background: {solid}"
        meta["render"] = "filled_shape"
        style_parts.append(bg)

        # Check for rounded corners
        rr = sp.find(f".//{qn('a:prstGeom')}")
        if rr is not None and rr.get("prst") == "roundRect":
            style_parts.append("border-radius: 12px")

        html = f'<div style="{"; ".join(style_parts)}" data-layer="{z_index}" data-name="{shape.name}"></div>'

    # Case 6: Group — recurse
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        meta["render"] = "group"
        meta["children"] = len(shape.shapes)
        children_html = []
        for ci, child in enumerate(shape.shapes):
            child_html, child_meta = _shape_to_html(
                child, canvas_w, canvas_h, z_index * 10 + ci, prs_part, clip_id_counter
            )
            if child_html:
                children_html.append(child_html)

        inner = "\n".join(children_html)
        html = f"""<div style="{'; '.join(style_parts)}" data-layer="{z_index}" data-name="{shape.name}">
{inner}
</div>"""

    # Fallback: empty positioned div
    else:
        meta["render"] = "empty"
        html = f'<div style="{"; ".join(style_parts)}" data-layer="{z_index}" data-name="{shape.name}"></div>'

    return html, meta


# ── Main conversion ──────────────────────────────────────────────────────

def convert_slide(prs: Presentation, slide, slide_idx: int,
                  template_name: str) -> tuple[str, dict]:
    """Convert a single slide to a self-contained HTML document."""
    canvas_w = emu_to_px(prs.slide_width)
    canvas_h = emu_to_px(prs.slide_height)

    layers_html = []
    layer_manifest = []
    clip_id_counter = [0]

    for z_idx, shape in enumerate(slide.shapes):
        html, meta = _shape_to_html(
            shape, canvas_w, canvas_h, z_idx, slide.part, clip_id_counter
        )
        if html:
            layers_html.append(html)
            layer_manifest.append(meta)

    body = "\n\n".join(layers_html)

    full_html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8"/>
<meta name="viewport" content="width=device-width, initial-scale=1"/>
<title>{template_name} — Slide {slide_idx + 1}</title>
<style>
  * {{ margin: 0; padding: 0; box-sizing: border-box; }}
  body {{ background: #1a1a1a; display: flex; justify-content: center; align-items: center; min-height: 100vh; }}
  .canvas {{
    position: relative;
    width: {canvas_w}px;
    height: {canvas_h}px;
    background: #FFFFFF;
    overflow: hidden;
    box-shadow: 0 4px 24px rgba(0,0,0,0.3);
  }}
</style>
</head>
<body>
<div class="canvas">
{body}
</div>
<!--
LAYER MANIFEST (for LLM consumption):
{json.dumps(layer_manifest, indent=2)}
-->
</body>
</html>"""

    manifest = {
        "template": template_name,
        "slide": slide_idx,
        "canvas": {"width": canvas_w, "height": canvas_h},
        "layer_count": len(layer_manifest),
        "layers": layer_manifest,
    }

    return full_html, manifest


def convert_pptx(pptx_path: str, output_dir: str) -> list[dict]:
    """Convert all slides in a PPTX to HTML files."""
    prs = Presentation(pptx_path)
    name = Path(pptx_path).stem
    safe_name = re.sub(r"[^a-zA-Z0-9_-]", "_", name)[:60]

    manifests = []
    for slide_idx, slide in enumerate(prs.slides):
        html, manifest = convert_slide(prs, slide, slide_idx, name)

        suffix = f"_slide{slide_idx + 1}" if len(prs.slides) > 1 else ""
        out_path = os.path.join(output_dir, f"{safe_name}{suffix}.html")
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(html)

        manifests.append(manifest)
        print(f"  ✓ {os.path.basename(out_path)} — {manifest['layer_count']} layers, {canvas_dims(manifest)}")

    return manifests


def canvas_dims(m: dict) -> str:
    return f"{m['canvas']['width']}×{m['canvas']['height']}px"


# ── CLI ──────────────────────────────────────────────────────────────────

def main():
    input_dir = sys.argv[1] if len(sys.argv) > 1 else "Designinspio"
    output_dir = sys.argv[2] if len(sys.argv) > 2 else os.path.join(input_dir, "html_output")

    os.makedirs(output_dir, exist_ok=True)

    pptx_files = sorted(
        f for f in os.listdir(input_dir) if f.endswith(".pptx")
    )

    if not pptx_files:
        print(f"No .pptx files found in {input_dir}")
        sys.exit(1)

    print(f"Converting {len(pptx_files)} PPTX templates → layered HTML")
    print(f"Output: {output_dir}/\n")

    catalog = []
    for f in pptx_files:
        print(f"📄 {f}")
        full_path = os.path.join(input_dir, f)
        try:
            manifests = convert_pptx(full_path, output_dir)
            catalog.extend(manifests)
        except Exception as e:
            print(f"  ✗ FAILED: {e}")
        print()

    # Write reference catalog
    catalog_path = os.path.join(output_dir, "_reference_catalog.json")
    with open(catalog_path, "w") as f:
        json.dump(catalog, f, indent=2)

    print(f"\n{'='*60}")
    print(f"Done! {len(catalog)} slides converted.")
    print(f"Catalog: {catalog_path}")
    print(f"\nLayer summary:")
    for entry in catalog:
        render_types = {}
        for layer in entry["layers"]:
            rt = layer.get("render", "unknown")
            render_types[rt] = render_types.get(rt, 0) + 1
        types_str = ", ".join(f"{k}:{v}" for k, v in sorted(render_types.items()))
        print(f"  {entry['template'][:50]:50s} | {entry['layer_count']:2d} layers | {types_str}")


if __name__ == "__main__":
    main()

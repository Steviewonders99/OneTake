#!/usr/bin/env python3
"""PPTX → LLM-ready reference code extractor.

Produces COMPACT HTML/CSS reference files that show layered design techniques
WITHOUT base64-embedded images. Images are replaced with colored placeholders
showing dimensions and clip-path shapes.

Output files are 2-15KB — small enough to inject into LLM prompts as
"here's how a professional designer builds layered creatives."

Usage:
    python scripts/pptx_to_reference_code.py [input_dir] [output_dir]
"""
from __future__ import annotations

import json
import math
import os
import re
import sys
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ── Constants ─────────────────────────────────────────────────────────────
EMU_PER_PX = 914400 / 96

# Placeholder colors by layer purpose (helps LLM understand what each layer IS)
LAYER_COLORS = [
    "#E8D5F5",  # z0 — light purple (background shapes)
    "#D5E8F5",  # z1 — light blue
    "#F5E8D5",  # z2 — light orange
    "#D5F5E8",  # z3 — light green
    "#F5D5E8",  # z4 — light pink
    "#E8F5D5",  # z5 — light lime
    "#D5D5F5",  # z6 — light indigo
    "#F5F5D5",  # z7 — light yellow
    "#E8E8E8",  # z8+ — light gray
]

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}


def emu_to_px(emu: int) -> float:
    return round(emu / EMU_PER_PX, 1)


def _parse_color(elem) -> str:
    for tag in ["a:srgbClr", "a:sysClr", "a:schemeClr"]:
        node = elem.find(qn(tag))
        if node is not None:
            if tag == "a:srgbClr":
                return f"#{node.get('val', '888888')}"
            elif tag == "a:sysClr":
                return f"#{node.get('lastClr', node.get('val', '888888'))}"
            else:
                scheme = {
                    "bg1": "#FFFFFF", "tx1": "#000000", "dk1": "#000000",
                    "lt1": "#FFFFFF", "accent1": "#0452BF", "accent2": "#CD128A",
                    "dk2": "#1A1A1A", "lt2": "#F5F5F5",
                }
                return scheme.get(node.get("val", ""), "#888888")
    return "#888888"


def _extract_gradient_css(sp) -> str | None:
    grad = sp.find(f".//{qn('a:gradFill')}")
    if grad is None:
        return None
    gs_lst = grad.find(qn("a:gsLst"))
    if gs_lst is None:
        return None

    stops = []
    for gs in gs_lst.findall(qn("a:gs")):
        pos = int(gs.get("pos", "0")) / 1000
        color = _parse_color(gs)
        alpha_elem = gs.find(f".//{qn('a:alpha')}")
        if alpha_elem is not None:
            alpha = int(alpha_elem.get("val", "100000")) / 100000
            r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
            stops.append(f"rgba({r},{g},{b},{alpha:.2f}) {pos:.0f}%")
        else:
            stops.append(f"{color} {pos:.0f}%")

    lin = grad.find(qn("a:lin"))
    angle = 180
    if lin is not None:
        angle = int(lin.get("ang", "0")) / 60000

    return f"linear-gradient({angle:.0f}deg, {', '.join(stops)})"


def _extract_solid_css(sp) -> str | None:
    solid = sp.find(f".//{qn('a:solidFill')}")
    if solid is None:
        return None
    color = _parse_color(solid)
    alpha_elem = solid.find(f".//{qn('a:alpha')}")
    if alpha_elem is not None:
        alpha = int(alpha_elem.get("val", "100000")) / 100000
        r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
        return f"rgba({r},{g},{b},{alpha:.2f})"
    return color


def _extract_svg_path(shape) -> dict | None:
    sp = shape._element
    cust_geom = sp.find(f".//{qn('a:custGeom')}")
    if cust_geom is None:
        return None
    path_lst = cust_geom.find(qn("a:pathLst"))
    if path_lst is None:
        return None

    for path_elem in path_lst.findall(qn("a:path")):
        pw = int(path_elem.get("w", "1"))
        ph = int(path_elem.get("h", "1"))
        d_parts = []
        for child in path_elem:
            tag = child.tag.split("}")[-1]
            if tag == "moveTo":
                pt = child.find(qn("a:pt"))
                if pt is not None:
                    d_parts.append(f"M{pt.get('x','0')},{pt.get('y','0')}")
            elif tag == "lnTo":
                pt = child.find(qn("a:pt"))
                if pt is not None:
                    d_parts.append(f"L{pt.get('x','0')},{pt.get('y','0')}")
            elif tag == "cubicBezTo":
                pts = child.findall(qn("a:pt"))
                if len(pts) == 3:
                    c = " ".join(f"{p.get('x','0')},{p.get('y','0')}" for p in pts)
                    d_parts.append(f"C{c}")
            elif tag == "close":
                d_parts.append("Z")
        if d_parts:
            return {"d": "".join(d_parts), "vb": f"0 0 {pw} {ph}"}
    return None


def _has_image_fill(shape) -> bool:
    sp = shape._element
    return (
        sp.find(f".//{qn('a:blipFill')}") is not None
        or sp.find(f".//{qn('p:blipFill')}") is not None
        or shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )


def _extract_text(shape) -> list[dict] | None:
    if not shape.has_text_frame or not shape.text.strip():
        return None
    runs = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            rd: dict = {"t": run.text}
            f = run.font
            if f.size:
                rd["sz"] = round(f.size.pt * 1.333, 1)
            if f.bold:
                rd["b"] = True
            if f.italic:
                rd["i"] = True
            try:
                if f.color and f.color.rgb:
                    rd["c"] = f"#{f.color.rgb}"
            except AttributeError:
                pass
            if f.name:
                rd["ff"] = f.name
            runs.append(rd)
        if para.runs:
            runs.append({"br": True})
    return runs if runs else None


# ── Shape → compact CSS/HTML ─────────────────────────────────────────────

def shape_to_code(shape, z: int, canvas_w: float, canvas_h: float,
                  clip_counter: list) -> tuple[str, dict]:
    """Convert shape to compact HTML + metadata dict."""
    left = emu_to_px(shape.left)
    top = emu_to_px(shape.top)
    w = emu_to_px(shape.width)
    h = emu_to_px(shape.height)
    rot = shape.rotation or 0
    sp = shape._element

    meta: dict = {
        "z": z,
        "name": shape.name,
        "pos": f"{left},{top} {w}×{h}",
    }
    if rot:
        meta["rot"] = f"{rot:.1f}°"

    # Classify the layer
    svg_path = _extract_svg_path(shape)
    has_img = _has_image_fill(shape)
    gradient = _extract_gradient_css(sp)
    solid = _extract_solid_css(sp)
    text_runs = _extract_text(shape)
    placeholder_color = LAYER_COLORS[min(z, len(LAYER_COLORS) - 1)]

    # Base style
    base = f"position:absolute;left:{left}px;top:{top}px;width:{w}px;height:{h}px;z-index:{z}"
    if rot:
        base += f";transform:rotate({rot:.1f}deg)"

    html = ""

    # ── Freeform + image = clipped photo ──
    if svg_path and has_img:
        cid = f"c{clip_counter[0]}"
        clip_counter[0] += 1
        meta["type"] = "clipped_photo"
        meta["technique"] = "SVG clip-path masks photo into organic blob shape"
        html = f'''<div style="{base}" data-layer="z{z}" data-role="clipped-photo">
  <svg viewBox="{svg_path['vb']}" width="100%" height="100%" style="position:absolute;inset:0">
    <defs><clipPath id="{cid}"><path d="{svg_path['d']}"/></clipPath></defs>
    <rect width="100%" height="100%" fill="{placeholder_color}" clip-path="url(#{cid})"/>
    <!-- PHOTO: Replace rect with <image href="ACTOR_PHOTO_URL" .../> -->
  </svg>
</div>'''

    # ── Freeform + gradient = decorative gradient blob ──
    elif svg_path and gradient:
        meta["type"] = "gradient_blob"
        meta["technique"] = "Organic SVG shape filled with brand gradient"
        meta["gradient"] = gradient
        gid = f"g{clip_counter[0]}"
        clip_counter[0] += 1
        # Extract first/last colors from gradient string
        colors = re.findall(r'(#[0-9a-fA-F]{6}|rgba?\([^)]+\))', gradient)
        c1 = colors[0] if colors else "#0452BF"
        c2 = colors[-1] if len(colors) > 1 else "#CD128A"
        html = f'''<div style="{base}" data-layer="z{z}" data-role="gradient-blob">
  <svg viewBox="{svg_path['vb']}" width="100%" height="100%" style="position:absolute;inset:0">
    <defs><linearGradient id="{gid}"><stop offset="0%" stop-color="{c1}"/><stop offset="100%" stop-color="{c2}"/></linearGradient></defs>
    <path d="{svg_path['d']}" fill="url(#{gid})"/>
  </svg>
</div>'''

    # ── Freeform + solid = decorative solid shape ──
    elif svg_path and solid:
        meta["type"] = "solid_shape"
        meta["technique"] = "Organic SVG shape with solid fill"
        meta["fill"] = solid
        html = f'''<div style="{base}" data-layer="z{z}" data-role="solid-shape">
  <svg viewBox="{svg_path['vb']}" width="100%" height="100%" style="position:absolute;inset:0">
    <path d="{svg_path['d']}" fill="{solid}"/>
  </svg>
</div>'''

    # ── Freeform with no fill = structural shape ──
    elif svg_path:
        meta["type"] = "structural_shape"
        meta["technique"] = "SVG shape used for layout structure"
        html = f'''<div style="{base}" data-layer="z{z}" data-role="structure">
  <svg viewBox="{svg_path['vb']}" width="100%" height="100%" style="position:absolute;inset:0">
    <path d="{svg_path['d']}" fill="{placeholder_color}" opacity="0.15"/>
  </svg>
</div>'''

    # ── Text box ──
    elif text_runs:
        meta["type"] = "text"
        # Build compact text
        text_parts = []
        for r in text_runs:
            if r.get("br"):
                text_parts.append("<br/>")
                continue
            styles = []
            if "sz" in r:
                styles.append(f"font-size:{r['sz']}px")
            if r.get("b"):
                styles.append("font-weight:700")
            if r.get("i"):
                styles.append("font-style:italic")
            if "c" in r:
                styles.append(f"color:{r['c']}")
            if "ff" in r:
                styles.append(f"font-family:'{r['ff']}',system-ui")
            s = f' style="{";".join(styles)}"' if styles else ""
            text_parts.append(f"<span{s}>{r['t']}</span>")
        inner = "".join(text_parts)
        meta["text"] = "".join(r.get("t", "") for r in text_runs if "t" in r)[:80]

        # Alignment
        align = "left"
        if shape.has_text_frame and shape.text_frame.paragraphs:
            al = shape.text_frame.paragraphs[0].alignment
            if al is not None:
                align = {0: "left", 1: "center", 2: "right", 3: "justify"}.get(al, "left")

        html = f'''<div style="{base};display:flex;align-items:center;text-align:{align};line-height:1.15" data-layer="z{z}" data-role="text">
  <div style="width:100%">{inner}</div>
</div>'''

    # ── Image without freeform path ──
    elif has_img:
        meta["type"] = "photo"
        meta["technique"] = "Rectangular photo placement"
        html = f'''<div style="{base};background:{placeholder_color}" data-layer="z{z}" data-role="photo">
  <!-- PHOTO: {w}×{h}px — object-fit:cover -->
</div>'''

    # ── Rectangle/shape with gradient ──
    elif gradient:
        meta["type"] = "gradient_rect"
        meta["gradient"] = gradient
        html = f'<div style="{base};background:{gradient}" data-layer="z{z}" data-role="gradient-bg"></div>'

    # ── Rectangle/shape with solid fill ──
    elif solid:
        meta["type"] = "filled_rect"
        meta["fill"] = solid
        rr = sp.find(f".//{qn('a:prstGeom')}")
        radius = ";border-radius:12px" if (rr is not None and rr.get("prst") == "roundRect") else ""
        html = f'<div style="{base};background:{solid}{radius}" data-layer="z{z}" data-role="filled-rect"></div>'

    # ── Group ──
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        meta["type"] = "group"
        meta["children"] = len(shape.shapes)
        children = []
        for ci, child in enumerate(shape.shapes):
            ch, cm = shape_to_code(child, z * 10 + ci, canvas_w, canvas_h, clip_counter)
            if ch:
                children.append(ch)
        inner = "\n".join(children)
        html = f'''<div style="{base}" data-layer="z{z}" data-role="group">
{inner}
</div>'''

    else:
        meta["type"] = "empty"
        return "", meta

    return html, meta


def convert_slide_compact(prs: Presentation, slide, slide_idx: int,
                          template_name: str) -> tuple[str, dict]:
    """Convert slide to compact LLM-ready HTML."""
    cw = emu_to_px(prs.slide_width)
    ch = emu_to_px(prs.slide_height)

    layers = []
    manifest = []
    clip_counter = [0]

    for z, shape in enumerate(slide.shapes):
        html, meta = shape_to_code(shape, z, cw, ch, clip_counter)
        if html:
            layers.append(html)
            manifest.append(meta)

    # Build compact summary comment
    types_summary = {}
    for m in manifest:
        t = m.get("type", "unknown")
        types_summary[t] = types_summary.get(t, 0) + 1
    summary = ", ".join(f"{v}× {k}" for k, v in sorted(types_summary.items(), key=lambda x: -x[1]))

    body = "\n\n".join(layers)

    full = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8"/>
<title>REF: {template_name[:60]}</title>
<style>
* {{ margin:0; padding:0; box-sizing:border-box; }}
body {{ background:#111; display:flex; justify-content:center; align-items:center; min-height:100vh; }}
.canvas {{ position:relative; width:{cw}px; height:{ch}px; background:#FFF; overflow:hidden; }}
</style>
</head>
<body>
<!--
DESIGN REFERENCE: {template_name}
Canvas: {cw}×{ch}px
Layers: {len(manifest)} ({summary})
Technique: Layered absolute positioning with SVG clip-paths for organic shapes,
CSS gradients for brand colors, z-index stacking for depth.

LAYER MANIFEST:
{json.dumps(manifest, indent=1, ensure_ascii=False)}
-->
<div class="canvas">
{body}
</div>
</body>
</html>"""

    return full, {
        "template": template_name,
        "slide": slide_idx,
        "canvas": f"{cw}×{ch}",
        "layers": len(manifest),
        "types": types_summary,
        "manifest": manifest,
    }


def main():
    input_dir = sys.argv[1] if len(sys.argv) > 1 else "Designinspio"
    output_dir = sys.argv[2] if len(sys.argv) > 2 else os.path.join(input_dir, "reference_code")

    os.makedirs(output_dir, exist_ok=True)

    pptx_files = sorted(f for f in os.listdir(input_dir) if f.endswith(".pptx"))
    if not pptx_files:
        print(f"No .pptx files in {input_dir}")
        sys.exit(1)

    print(f"Extracting {len(pptx_files)} PPTX → LLM-ready reference code")
    print(f"Output: {output_dir}/\n")

    catalog = []
    for f in pptx_files:
        prs = Presentation(os.path.join(input_dir, f))
        name = Path(f).stem
        safe = re.sub(r"[^a-zA-Z0-9_-]", "_", name)[:60]

        for si, slide in enumerate(prs.slides):
            suffix = f"_s{si+1}" if len(prs.slides) > 1 else ""
            try:
                html, meta = convert_slide_compact(prs, slide, si, name)
                out = os.path.join(output_dir, f"{safe}{suffix}.html")
                with open(out, "w") as fh:
                    fh.write(html)

                size_kb = len(html.encode()) / 1024
                print(f"  ✓ {os.path.basename(out):60s} {meta['layers']:2d} layers  {size_kb:.1f}KB")
                catalog.append(meta)
            except Exception as e:
                print(f"  ✗ {safe}{suffix}: {e}")

    # Write compact catalog
    cat_path = os.path.join(output_dir, "_catalog.json")
    with open(cat_path, "w") as f:
        json.dump(catalog, f, indent=2, ensure_ascii=False)

    print(f"\nDone! {len(catalog)} references extracted → {output_dir}/")
    print(f"\nSize summary:")
    total = 0
    for fn in sorted(os.listdir(output_dir)):
        if not fn.endswith(".html"):
            continue
        fp = os.path.join(output_dir, fn)
        sz = os.path.getsize(fp)
        total += sz
        print(f"  {fn:65s} {sz/1024:6.1f}KB")
    print(f"\n  Total: {total/1024:.1f}KB ({total/1024/1024:.2f}MB)")


if __name__ == "__main__":
    main()
